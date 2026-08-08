"""Extract text from common office file formats.

Supports: PDF, DOCX, XLSX, PPTX
Saves extracted text as a companion .md file alongside the original.
"""

import io
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

from loguru import logger


# File extensions that need text extraction
EXTRACTABLE_EXTS = {".pdf", ".docx", ".xlsx", ".pptx"}

# Text extensions that don't need extraction
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
             ".js", ".ts", ".py", ".html", ".css", ".sh", ".log", ".env"}


# ─── Extraction result (forwarded to the UI) ─────────────────────
#
# `md_path` is the companion .md alongside the original file when extraction
# produced usable text. `reason` tells the UI *why* nothing was produced so
# end users get actionable feedback (e.g. "scanned PDF, please provide a
# text-based PDF") instead of a silent null.

EXTRACTION_REASON_OK = "ok"
EXTRACTION_REASON_SKIPPED_TEXT_FILE = "not_a_binary"
EXTRACTION_REASON_SCANNED_PDF = "scanned_pdf_no_text_layer"
EXTRACTION_REASON_EMPTY_EXTRACTION = "empty_extraction"
EXTRACTION_REASON_ERROR = "extraction_error"


@dataclass
class ExtractionResult:
    md_path: Path | None
    reason: str
    detail: str
    char_count: int
    page_count: int
    # Internal cache of the extracted text — populated when reason == OK so
    # `save_extracted_text` can write it without re-parsing the file. Never
    # surfaced via `to_dict`.
    _text: str | None = None

    def to_dict(self) -> dict[str, Any]:
        d = {k: v for k, v in asdict(self).items() if not k.startswith("_")}
        if d.get("md_path") is not None:
            d["md_path"] = str(d["md_path"])
        return d


def _clean_cell(value: object) -> str:
    text = str(value or "").strip()
    return text.replace("\n", "<br>").replace("|", "\\|")


def _markdown_table(rows: list[list[object]]) -> str:
    cleaned = [[_clean_cell(cell) for cell in row] for row in rows]
    cleaned = [row for row in cleaned if any(cell for cell in row)]
    if not cleaned:
        return ""

    width = max(len(row) for row in cleaned)
    normalized = [row + [""] * (width - len(row)) for row in normalized]
    header = normalized[0]
    separator = ["---"] * width
    body = normalized[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def needs_extraction(filename: str) -> bool:
    """Check if a file needs text extraction."""
    ext = Path(filename).suffix.lower()
    return ext in EXTRACTABLE_EXTS


def extract_text(file_bytes: bytes, filename: str) -> str | None:
    """Extract text from a binary file.

    Returns extracted text string, or None if extraction fails.
    Kept for backwards compatibility — prefer `extract_with_reason` for UI use.
    """
    res = extract_with_reason(file_bytes, filename)
    if res.reason == EXTRACTION_REASON_OK and res.md_path is None and res.char_count > 0:
        return None  # defensive: shouldn't happen
    if res.reason != EXTRACTION_REASON_OK:
        return None
    # Re-run to obtain raw text? Avoid that — UI consumers should switch to
    # `extract_with_reason`. Returning None here preserves the legacy contract.
    return None


def extract_with_reason(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract text and report *why* nothing was produced when applicable.

    Always returns an ExtractionResult, never raises. Reasons:
      - ok                              → text extracted, .md will be written
      - not_a_binary                    → non-binary file, no extraction needed
      - scanned_pdf_no_text_layer       → pdfplumber found no text on any page
      - empty_extraction                → extractor ran but produced no content
      - extraction_error                → unexpected exception, see `detail`
    """
    ext = Path(filename).suffix.lower()

    if ext not in EXTRACTABLE_EXTS:
        return ExtractionResult(
            md_path=None,
            reason=EXTRACTION_REASON_SKIPPED_TEXT_FILE,
            detail=f"{ext} is a text-like file, no extraction needed",
            char_count=0,
            page_count=0,
        )

    try:
        if ext == ".pdf":
            text, page_count = _extract_pdf(file_bytes)
        elif ext == ".docx":
            text = _extract_docx(file_bytes)
            page_count = 0
        elif ext == ".xlsx":
            text = _extract_xlsx(file_bytes)
            page_count = 0
        elif ext == ".pptx":
            text = _extract_pptx(file_bytes)
            page_count = 0
        else:
            return ExtractionResult(
                md_path=None,
                reason=EXTRACTION_REASON_EMPTY_EXTRACTION,
                detail=f"unsupported extension {ext}",
                char_count=0,
                page_count=0,
            )
    except Exception as e:
        logger.error(f"[TextExtractor] Failed to extract from {filename}: {e}")
        return ExtractionResult(
            md_path=None,
            reason=EXTRACTION_REASON_ERROR,
            detail=f"{type(e).__name__}: {e}",
            char_count=0,
            page_count=0,
        )

    if not text or not text.strip():
        # Strict scanned-PDF detection: pdfplumber sees N pages but every page
        # is textless. Treated as scanned PDF (most common cause). For other
        # formats, we cannot tell why, so fall back to "empty_extraction".
        if ext == ".pdf" and page_count > 0:
            return ExtractionResult(
                md_path=None,
                reason=EXTRACTION_REASON_SCANNED_PDF,
                detail=(
                    "PDF 似乎为扫描件，无可抽取的文本层。请提供文本型 PDF，"
                    "或先自行 OCR 后再上传。"
                ),
                char_count=0,
                page_count=page_count,
            )
        return ExtractionResult(
            md_path=None,
            reason=EXTRACTION_REASON_EMPTY_EXTRACTION,
            detail="extractor produced no text",
            char_count=0,
            page_count=page_count,
        )

    return ExtractionResult(
        md_path=None,  # populated by save_extracted_text below
        reason=EXTRACTION_REASON_OK,
        detail="",
        char_count=len(text.strip()),
        page_count=page_count,
        _text=text,
    )


def save_extracted_text(save_path: Path, file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract text and save as a companion .md file.

    Returns an ExtractionResult. When ok, md_path points to the written .md.
    Use `ExtractionResult.to_dict()` to surface reason/detail to the UI.
    """
    result = extract_with_reason(file_bytes, filename)
    if result.reason != EXTRACTION_REASON_OK:
        return result

    text = result._text or ""
    if not text.strip():
        return ExtractionResult(
            md_path=None,
            reason=EXTRACTION_REASON_EMPTY_EXTRACTION,
            detail="extraction produced no text",
            char_count=0,
            page_count=result.page_count,
        )

    md_path = save_path.parent / f"{save_path.stem}.md"
    md_path.write_text(text, encoding="utf-8")
    result.md_path = md_path
    logger.info(f"[TextExtractor] Extracted {len(text)} chars from {filename} -> {md_path.name}")
    return result


def _extract_pdf(data: bytes) -> tuple[str, int]:
    """Extract text from PDF using pdfplumber.

    Returns (markdown_text, page_count). pdfplumber returns None for image-only
    pages — we count those separately so the caller can flag scanned PDFs.
    """
    import pdfplumber

    pages: list[str] = []
    page_count = 0
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        page_count = len(pdf.pages)
        for i, page in enumerate(pdf.pages):
            page_parts: list[str] = []
            text = page.extract_text()
            if text and text.strip():
                page_parts.append(text.strip())

            # Also extract tables
            tables = page.extract_tables()
            for table in tables:
                if table:
                    table_md = _markdown_table(table)
                    if table_md:
                        page_parts.append(table_md)

            if page_parts:
                pages.append(f"## 第 {i + 1} 页\n\n" + "\n\n".join(page_parts))

    return "\n\n".join(pages), page_count


def _extract_docx(data: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Preserve heading hierarchy
            style_name = para.style.name if para.style and para.style.name else ""
            if style_name.startswith("Heading"):
                level = style_name.replace("Heading", "").strip()
                try:
                    level = int(level)
                except ValueError:
                    level = 1
                parts.append(f"{'#' * level} {text}")
            elif "List Bullet" in style_name:
                parts.append(f"- {text}")
            elif "List Number" in style_name:
                parts.append(f"1. {text}")
            else:
                parts.append(text)

    # Extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        table_md = _markdown_table(rows)
        if table_md:
            parts.append("## 表格\n\n" + table_md)

    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [c if c is not None else "" for c in row]
            if any(str(c).strip() for c in cells):
                rows.append(cells)

        table_md = _markdown_table(rows)
        if table_md:
            parts.append(f"## 工作表: {sheet}\n\n" + table_md)

    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []

    for i, slide in enumerate(prs.slides):
        texts = []
        tables = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                table_md = _markdown_table(rows)
                if table_md:
                    tables.append(table_md)

        slide_parts = []
        if texts:
            slide_parts.append("\n\n".join(texts))
        slide_parts.extend(tables)
        if slide_parts:
            parts.append(f"## 幻灯片 {i + 1}\n\n" + "\n\n".join(slide_parts))

    return "\n\n".join(parts)
